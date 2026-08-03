from __future__ import annotations

import csv
import json
import re
import zipfile
from collections import Counter, defaultdict
from pathlib import Path
from xml.etree import ElementTree as ET

import pdfplumber
import yaml
from docx import Document
from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
LESSONS = ROOT / "lessons"
BUILD = ROOT / "build" / "7b_acceptance"
REPORTS = ROOT / "reports"
EXPECTED = {
    "教学设计.docx": ("教学设计", ".docx", 4),
    "教学设计.pdf": ("教学设计", ".pdf", 4),
    "课堂教学.pptx": ("PPT", ".pptx", 15),
    "学生学案.docx": ("学案", ".docx", 3),
    "学生学案.pdf": ("学案", ".pdf", 3),
    "学案教师版.docx": ("学案", ".docx", 3),
    "学案教师版.pdf": ("学案", ".pdf", 3),
}


def compact(text: str) -> str:
    return re.sub(r"\s+", "", text or "")


def docx_text(path: Path) -> str:
    doc = Document(path)
    chunks = [p.text for p in doc.paragraphs]
    for table in doc.tables:
        for row in table.rows:
            chunks.extend(cell.text for cell in row.cells)
    return "\n".join(chunks)


def docx_pages(path: Path) -> int:
    with zipfile.ZipFile(path) as archive:
        xml = archive.read("docProps/app.xml")
    root = ET.fromstring(xml)
    for node in root.iter():
        if node.tag.endswith("Pages"):
            return int(node.text or 0)
    return 0


def pdf_details(path: Path) -> dict:
    with pdfplumber.open(path) as pdf:
        texts = [page.extract_text() or "" for page in pdf.pages]
        sizes = [(round(page.width, 2), round(page.height, 2)) for page in pdf.pages]
        metadata = pdf.metadata or {}
    return {
        "pages": len(texts),
        "texts": texts,
        "text": "\n".join(texts),
        "sizes": sizes,
        "producer": metadata.get("Producer", ""),
    }


def ppt_details(path: Path) -> dict:
    prs = Presentation(path)
    texts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                texts.append(shape.text)
    with zipfile.ZipFile(path) as archive:
        notes = len([name for name in archive.namelist() if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", name)])
    return {
        "slides": len(prs.slides),
        "width": round(prs.slide_width / 914400, 2),
        "height": round(prs.slide_height / 914400, 2),
        "notes": notes,
        "text": "\n".join(texts),
    }


def locate_sources() -> list[tuple[Path, dict]]:
    sources: list[tuple[Path, dict]] = []
    for path in sorted(LESSONS.rglob("lesson.yml")):
        data = yaml.safe_load(path.read_text(encoding="utf-8"))
        if str(data.get("meta", {}).get("lesson_id", "")).startswith("C0") or str(data.get("meta", {}).get("lesson_id", "")).startswith("C1"):
            lesson_dir = path.parent.parent
            sources.append((lesson_dir, data))
    return sources


def add(checks: list[dict], lesson_id: str, category: str, check: str, passed: bool, detail: str = "") -> None:
    checks.append({
        "lesson_id": lesson_id,
        "category": category,
        "check": check,
        "status": "PASS" if passed else "FAIL",
        "detail": detail,
    })


def main() -> None:
    BUILD.mkdir(parents=True, exist_ok=True)
    REPORTS.mkdir(parents=True, exist_ok=True)
    sources = locate_sources()
    checks: list[dict] = []
    inventory: list[dict] = []
    chapter_summary: dict[str, Counter] = defaultdict(Counter)

    add(checks, "ALL", "structure", "lesson_source_count", len(sources) == 59, str(len(sources)))
    ids = [data["meta"]["lesson_id"] for _, data in sources]
    add(checks, "ALL", "structure", "unique_lesson_ids", len(ids) == len(set(ids)) == 59, str(len(set(ids))))

    ppt_audit = {}
    audit_csv = ROOT / "build" / "7b_batch" / "ppt_audit_final.csv"
    with audit_csv.open(encoding="utf-8-sig", newline="") as handle:
        for row in csv.DictReader(handle):
            ppt_audit[str((ROOT / row["file"]).resolve()).lower()] = row

    live_word_pages: dict[str, int] = {}
    live_page_csv = BUILD / "word_live_page_check.csv"
    with live_page_csv.open(encoding="utf-8-sig", newline="") as handle:
        for row in csv.DictReader(handle):
            live_word_pages[str(Path(row["Docx"]).resolve()).lower()] = int(row["Pages"])

    for lesson_dir, data in sources:
        meta = data["meta"]
        lesson_id = meta["lesson_id"]
        chapter = lesson_id[:3]
        prefix = meta["file_prefix"]
        chapter_summary[chapter]["lessons"] += 1
        expected_paths = {
            "教学设计.docx": lesson_dir / "教学设计" / f"{prefix}_教学设计.docx",
            "教学设计.pdf": lesson_dir / "教学设计" / f"{prefix}_教学设计.pdf",
            "课堂教学.pptx": lesson_dir / "PPT" / f"{prefix}_课堂教学.pptx",
            "学生学案.docx": lesson_dir / "学案" / f"{prefix}_学生学案.docx",
            "学生学案.pdf": lesson_dir / "学案" / f"{prefix}_学生学案.pdf",
            "学案教师版.docx": lesson_dir / "学案" / f"{prefix}_学案教师版.docx",
            "学案教师版.pdf": lesson_dir / "学案" / f"{prefix}_学案教师版.pdf",
        }
        for key, path in expected_paths.items():
            exists = path.exists() and path.stat().st_size > 1000
            add(checks, lesson_id, "files", key, exists, str(path.relative_to(ROOT)))
            if exists:
                chapter_summary[chapter][key] += 1
                inventory.append({
                    "lesson_id": lesson_id,
                    "chapter": chapter,
                    "title": meta["title"],
                    "type": key,
                    "path": str(path.relative_to(ROOT)),
                    "size_bytes": path.stat().st_size,
                })
        if not all(path.exists() for path in expected_paths.values()):
            continue

        title = compact(meta["title"])
        teaching_path = expected_paths["教学设计.docx"]
        student_path = expected_paths["学生学案.docx"]
        teacher_path = expected_paths["学案教师版.docx"]
        teaching_text = docx_text(teaching_path)
        student_text = docx_text(student_path)
        teacher_text = docx_text(teacher_path)

        for label, path, text in (
            ("teaching", teaching_path, teaching_text),
            ("student", student_path, student_text),
            ("teacher", teacher_path, teacher_text),
        ):
            add(checks, lesson_id, "docx", f"{label}_opens", bool(text.strip()), f"chars={len(text)}")
            add(checks, lesson_id, "docx", f"{label}_title", title in compact(text), meta["title"])
            add(checks, lesson_id, "docx", f"{label}_no_placeholders", "[[" not in text and "]]" not in text, "")

        leaks = [token for token in ("参考答案", "评分点", "易错点") if token in student_text]
        add(checks, lesson_id, "content", "student_answer_isolation", not leaks, ",".join(leaks))
        teacher_tokens = [token for token in ("参考答案", "评分点", "易错点") if token in teacher_text]
        add(checks, lesson_id, "content", "teacher_guidance_complete", len(teacher_tokens) == 3, ",".join(teacher_tokens))

        questions = data.get("questions", [])
        add(checks, lesson_id, "source", "question_count", len(questions) == 8, str(len(questions)))
        for index, question in enumerate(questions, 1):
            prompt = compact(str(question.get("prompt", "")))
            answer = compact(str(question.get("answer", "")))
            add(checks, lesson_id, "content", f"Q{index}_student_prompt", bool(prompt) and prompt in compact(student_text), prompt[:40])
            add(checks, lesson_id, "content", f"Q{index}_teacher_prompt", bool(prompt) and prompt in compact(teacher_text), prompt[:40])
            add(checks, lesson_id, "content", f"Q{index}_teacher_answer", bool(answer) and answer in compact(teacher_text), answer[:40])

        for label, docx_path, pdf_path, expected_pages in (
            ("teaching", teaching_path, expected_paths["教学设计.pdf"], 4),
            ("student", student_path, expected_paths["学生学案.pdf"], 3),
            ("teacher", teacher_path, expected_paths["学案教师版.pdf"], 3),
        ):
            details = pdf_details(pdf_path)
            live_pages = live_word_pages.get(str(docx_path.resolve()).lower(), 0)
            add(checks, lesson_id, "pdf", f"{label}_page_count", details["pages"] == expected_pages, f"pdf={details['pages']}; expected={expected_pages}")
            add(checks, lesson_id, "pdf", f"{label}_word_pdf_pages", live_pages == details["pages"], f"word={live_pages}; pdf={details['pages']}")
            add(checks, lesson_id, "pdf", f"{label}_a4", all(abs(w - 595.32) < 1 and abs(h - 841.92) < 1 for w, h in details["sizes"]), str(details["sizes"][:1]))
            # User requirements prefer Microsoft Word, but explicitly allow a
            # headless LibreOffice export when Word automation is unavailable.
            producer = details["producer"]
            direct_office_export = ("Microsoft" in producer and "Word" in producer) or "LibreOffice" in producer
            add(checks, lesson_id, "pdf", f"{label}_word_export", direct_office_export, producer)
            add(checks, lesson_id, "pdf", f"{label}_searchable", all(len(compact(text)) > 30 for text in details["texts"]), str([len(compact(t)) for t in details["texts"]]))
            add(checks, lesson_id, "pdf", f"{label}_title_consistent", title in compact(details["text"]), meta["title"])
            if label == "student":
                pdf_leaks = [token for token in ("参考答案", "评分点", "易错点") if token in details["text"]]
                add(checks, lesson_id, "pdf", "student_pdf_answer_isolation", not pdf_leaks, ",".join(pdf_leaks))
            if label == "teacher":
                add(checks, lesson_id, "pdf", "teacher_pdf_guidance_complete", all(token in details["text"] for token in ("参考答案", "评分点", "易错点")), "")

        ppt_path = expected_paths["课堂教学.pptx"]
        ppt = ppt_details(ppt_path)
        add(checks, lesson_id, "ppt", "opens_and_slide_count", ppt["slides"] == 15, str(ppt["slides"]))
        add(checks, lesson_id, "ppt", "ratio_16_9", abs(ppt["width"] / ppt["height"] - 16 / 9) < 0.01, f"{ppt['width']}x{ppt['height']}")
        add(checks, lesson_id, "ppt", "speaker_notes", ppt["notes"] == ppt["slides"], f"{ppt['notes']}/{ppt['slides']}")
        add(checks, lesson_id, "ppt", "title_consistent", title in compact(ppt["text"]), meta["title"])
        add(checks, lesson_id, "ppt", "no_placeholders", "[[" not in ppt["text"] and "]]" not in ppt["text"], "")
        audit = ppt_audit.get(str(ppt_path.resolve()).lower())
        add(checks, lesson_id, "ppt", "layout_audit", bool(audit) and int(audit["issue_count"]) == 0, "missing" if not audit else audit["issue_count"])

    rendered_word = ROOT / "build" / "7b_rendered_word"
    add(checks, "ALL", "render", "word_pdf_contact_sheets", len(list(rendered_word.rglob("contact_sheet.png"))) == 177, str(len(list(rendered_word.rglob("contact_sheet.png")))))
    add(checks, "ALL", "render", "word_pdf_page_images", len(list(rendered_word.rglob("page-*.png"))) == 590, str(len(list(rendered_word.rglob("page-*.png")))))
    rendered_ppt = ROOT / "build" / "7b_rendered_ppt"
    add(checks, "ALL", "render", "ppt_contact_sheets", len(list(rendered_ppt.rglob("contact_sheet.png"))) == 59, str(len(list(rendered_ppt.rglob("contact_sheet.png")))))
    add(checks, "ALL", "render", "ppt_slide_images", len(list(rendered_ppt.rglob("slide-*.png"))) == 885, str(len(list(rendered_ppt.rglob("slide-*.png")))))

    status_counts = Counter(row["status"] for row in checks)
    output_json = {
        "summary": {
            "lessons": len(sources),
            "core_files": len(inventory),
            "checks": len(checks),
            "passed": status_counts["PASS"],
            "failed": status_counts["FAIL"],
        },
        "chapter_summary": {key: dict(value) for key, value in chapter_summary.items()},
        "checks": checks,
        "inventory": inventory,
    }
    (BUILD / "full_validation.json").write_text(json.dumps(output_json, ensure_ascii=False, indent=2), encoding="utf-8")
    for filename, rows in (("full_validation.csv", checks), ("core_inventory.csv", inventory)):
        with (BUILD / filename).open("w", encoding="utf-8-sig", newline="") as handle:
            writer = csv.DictWriter(handle, fieldnames=list(rows[0]))
            writer.writeheader()
            writer.writerows(rows)

    failures = [row for row in checks if row["status"] == "FAIL"]
    lines = [
        "# 七年级下册全册自动质量检查",
        "",
        f"- 课时：{len(sources)}",
        f"- 核心成品：{len(inventory)}",
        f"- 检查项：{len(checks)}",
        f"- 通过：{status_counts['PASS']}",
        f"- 失败：{status_counts['FAIL']}",
        "",
        "## 失败项",
        "",
    ]
    if failures:
        lines.extend(["| 课时 | 类别 | 检查 | 说明 |", "|---|---|---|---|"])
        lines.extend(f"| {r['lesson_id']} | {r['category']} | {r['check']} | {r['detail']} |" for r in failures)
    else:
        lines.append("无。")
    (REPORTS / "七下全册自动质量检查报告.md").write_text("\n".join(lines) + "\n", encoding="utf-8")
    print(json.dumps(output_json["summary"], ensure_ascii=False))
    if failures:
        print(f"failures={len(failures)}; see {REPORTS / '七下全册自动质量检查报告.md'}")
        raise SystemExit(1)


if __name__ == "__main__":
    main()
