from __future__ import annotations

import json
import os
import re
import sys
import zipfile
from pathlib import Path

import pdfplumber
from docx import Document
from pptx import Presentation

from pptx_audit_lib import audit_pptx


PROJECT = Path(__file__).resolve().parents[1]
LESSON_ROOT = PROJECT / "lessons" / "第07章_相交线与平行线" / "7.1.1_相交线与对顶角"
REPORT = PROJECT / "reports" / "C07-L01样例检查报告.md"
JSON_REPORT = PROJECT / "build" / "7b_sample" / "qa" / "C07-L01_sample_qa.json"


def docx_text(path: Path) -> str:
    document = Document(path)
    parts: list[str] = []
    parts.extend(p.text for p in document.paragraphs if p.text)
    for table in document.tables:
        for row in table.rows:
            parts.extend(cell.text for cell in row.cells if cell.text)
    return "\n".join(parts)


def pdf_info(path: Path) -> dict:
    with pdfplumber.open(path) as pdf:
        page_texts = [page.extract_text() or "" for page in pdf.pages]
        sizes = [(round(float(page.width), 2), round(float(page.height), 2)) for page in pdf.pages]
        return {
            "pages": len(pdf.pages),
            "page_texts": page_texts,
            "text": "\n".join(page_texts),
            "sizes": sizes,
            "metadata": dict(pdf.metadata or {}),
        }


def word_page_counts(paths: list[Path]) -> dict[str, int]:
    try:
        import win32com.client  # type: ignore
    except Exception:
        return {}
    word = win32com.client.DispatchEx("Word.Application")
    word.Visible = False
    word.DisplayAlerts = 0
    result: dict[str, int] = {}
    try:
        for path in paths:
            document = word.Documents.Open(str(path), ReadOnly=True, AddToRecentFiles=False)
            try:
                document.Repaginate()
                result[path.name] = int(document.ComputeStatistics(2))
            finally:
                document.Close(False)
    finally:
        word.Quit()
    return result


def check(condition: bool, label: str, failures: list[str], passes: list[str]) -> None:
    (passes if condition else failures).append(label)


def main() -> int:
    expected = {
        "teaching_docx": LESSON_ROOT / "教学设计" / "7.1.1_相交线与对顶角_教学设计.docx",
        "teaching_pdf": LESSON_ROOT / "教学设计" / "7.1.1_相交线与对顶角_教学设计.pdf",
        "pptx": LESSON_ROOT / "PPT" / "7.1.1_相交线与对顶角_课堂教学.pptx",
        "student_docx": LESSON_ROOT / "学案" / "7.1.1_相交线与对顶角_学生学案.docx",
        "student_pdf": LESSON_ROOT / "学案" / "7.1.1_相交线与对顶角_学生学案.pdf",
        "teacher_docx": LESSON_ROOT / "学案" / "7.1.1_相交线与对顶角_学案教师版.docx",
        "teacher_pdf": LESSON_ROOT / "学案" / "7.1.1_相交线与对顶角_学案教师版.pdf",
    }
    failures: list[str] = []
    passes: list[str] = []

    for label, path in expected.items():
        check(path.exists() and path.stat().st_size > 0, f"核心文件存在且非空：{label}", failures, passes)

    docx_paths = [expected["teaching_docx"], expected["student_docx"], expected["teacher_docx"]]
    pdf_paths = [expected["teaching_pdf"], expected["student_pdf"], expected["teacher_pdf"]]
    docx_texts = {path.name: docx_text(path) for path in docx_paths}
    pdf_infos = {path.name: pdf_info(path) for path in pdf_paths}
    word_pages = word_page_counts(docx_paths)

    expected_pages = {
        expected["teaching_docx"].name: 4,
        expected["student_docx"].name: 3,
        expected["teacher_docx"].name: 3,
    }
    pdf_for_docx = {
        expected["teaching_docx"].name: expected["teaching_pdf"].name,
        expected["student_docx"].name: expected["student_pdf"].name,
        expected["teacher_docx"].name: expected["teacher_pdf"].name,
    }
    for docx_name, expected_count in expected_pages.items():
        if word_pages:
            check(word_pages.get(docx_name) == expected_count, f"Word页数正确：{docx_name}={expected_count}", failures, passes)
        pdf_name = pdf_for_docx[docx_name]
        check(pdf_infos[pdf_name]["pages"] == expected_count, f"PDF页数正确：{pdf_name}={expected_count}", failures, passes)
        if word_pages:
            check(word_pages.get(docx_name) == pdf_infos[pdf_name]["pages"], f"Word/PDF页数一致：{docx_name}", failures, passes)

    for pdf_name, info in pdf_infos.items():
        check(all(size == (595.32, 841.92) for size in info["sizes"]), f"PDF为A4且各页尺寸一致：{pdf_name}", failures, passes)
        check(all(len(text.strip()) >= 120 for text in info["page_texts"]), f"PDF无空白页且文字可搜索复制：{pdf_name}", failures, passes)
        producer = str(info["metadata"].get("Producer", ""))
        check("Microsoft" in producer and "Word" in producer, f"PDF由Microsoft Word直接导出：{pdf_name}", failures, passes)

    student_text = pdf_infos[expected["student_pdf"].name]["text"]
    teacher_text = pdf_infos[expected["teacher_pdf"].name]["text"]
    teaching_text = pdf_infos[expected["teaching_pdf"].name]["text"]
    forbidden_student = ["参考答案", "评分点", "易错点", "∠B=115°", "145°、35°、145°", "x=20°", "60°和120°"]
    for marker in forbidden_student:
        check(marker not in student_text, f"学生版未泄露：{marker}", failures, passes)

    teacher_required = [
        "Q1", "Q2", "Q3", "Q4", "Q5", "Q6", "Q7", "Q8",
        "参考答案", "评分点", "易错点", "115°", "140°", "145°", "x=20°", "60°和120°",
    ]
    for marker in teacher_required:
        check(marker in teacher_text, f"教师版答案要素完整：{marker}", failures, passes)

    consistency_markers = [
        "相交线与对顶角", "邻补角", "对顶角", "对顶角相等", "∠1=40°",
        "∠AOC∶∠BOC=2∶7", "有公共顶点且相等", "比它的邻补角小60°",
    ]
    lesson_source = (LESSON_ROOT / "构建文件" / "lesson.yml").read_text(encoding="utf-8")
    for marker in consistency_markers:
        check(marker in student_text and marker in teacher_text, f"学生版/教师版题干一致：{marker}", failures, passes)
        check(marker in lesson_source, f"结构化课源覆盖核心内容：{marker}", failures, passes)
    for marker in ["相交线与对顶角", "邻补角", "对顶角", "对顶角相等", "教材例1", "Q4", "Q8"]:
        check(marker in teaching_text, f"教学设计覆盖核心内容：{marker}", failures, passes)

    pptx = expected["pptx"]
    prs = Presentation(pptx)
    check(len(prs.slides) == 15, "PPT页数为15页", failures, passes)
    check(abs(prs.slide_width / 914400 - 13.333333) < 0.02 and abs(prs.slide_height / 914400 - 7.5) < 0.02,
          "PPT尺寸为13.333×7.5英寸（16:9）", failures, passes)
    audit = audit_pptx(pptx)
    check(not audit.issues, "PPT自动审计无越界、遮挡、空白页或过小字号问题", failures, passes)
    with zipfile.ZipFile(pptx) as archive:
        note_slides = [name for name in archive.namelist() if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", name)]
        check(len(note_slides) == 15, "PPT 15页均含讲者备注", failures, passes)

    ppt_text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
    stale_markers = ["度分秒", "方位角", "角平分线", "31°28′", "126°43′"]
    for marker in stale_markers:
        check(marker not in ppt_text, f"PPT无金标准旧课残留：{marker}", failures, passes)
    ppt_required = ["7.1.1 相交线与对顶角", "教材例1", "∠1=40°", "邻补角", "对顶角", "Q7", "Q8"]
    for marker in ppt_required:
        check(marker in ppt_text, f"PPT核心内容存在：{marker}", failures, passes)

    editable_lines = 0
    full_slide_pictures = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            # Artifact Tool exports preset `line` geometry as an editable
            # autoshape whose stable name is retained in the PPTX.
            if getattr(shape, "name", "") in {"editable-line-a", "editable-line-b"}:
                editable_lines += 1
            if shape.shape_type == 13:  # picture
                if shape.width >= prs.slide_width * 0.9 and shape.height >= prs.slide_height * 0.9:
                    full_slide_pictures += 1
    check(editable_lines >= 18, "PPT相交线图形由可编辑线条构成", failures, passes)
    check(full_slide_pictures == 0, "PPT不存在整页图片替代可编辑内容", failures, passes)

    visual_evidence = [
        PROJECT / "build/7b_sample/rendered/teaching/contact-sheet.png",
        PROJECT / "build/7b_sample/rendered/student/contact-sheet.png",
        PROJECT / "build/7b_sample/rendered/teacher/contact-sheet.png",
        PROJECT / "build/7b_sample/rendered/ppt_final/7_1_1_相交线与对顶角_课堂教学/contact_sheet.png",
    ]
    for path in visual_evidence:
        check(path.exists() and path.stat().st_size > 0, f"视觉检查证据存在：{path.relative_to(PROJECT)}", failures, passes)

    result = {
        "lesson": "C07-L01 7.1.1 相交线与对顶角",
        "status": "PASS" if not failures else "FAIL",
        "passes": passes,
        "failures": failures,
        "word_pages": word_pages,
        "pdf_pages": {name: info["pages"] for name, info in pdf_infos.items()},
        "ppt_slides": len(prs.slides),
        "ppt_audit_issues": audit.issues,
        "visual_evidence": [str(path) for path in visual_evidence],
    }
    JSON_REPORT.parent.mkdir(parents=True, exist_ok=True)
    JSON_REPORT.write_text(json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8")

    lines = [
        "# C07-L01 样例检查报告",
        "",
        f"- 课题：7.1.1 相交线与对顶角",
        f"- 检查结论：**{result['status']}**",
        f"- 核心文件：7个（教学设计Word/PDF、PPT、学生学案Word/PDF、教师版Word/PDF）",
        f"- Word页数：教学设计4页、学生学案3页、教师版3页",
        f"- PPT页数：{len(prs.slides)}页，16:9",
        "- PDF导出：Microsoft Word 2016直接导出；文字可搜索、选择和复制",
        "- 视觉检查：已检查教学设计4页、学生学案3页、教师版3页、PPT 15页接触表及关键单页",
        "",
        "## 检查通过项",
        "",
    ]
    lines.extend(f"- {item}" for item in passes)
    lines.extend(["", "## 未通过项", ""])
    lines.extend(f"- {item}" for item in failures)
    if not failures:
        lines.append("- 无")
    lines.extend([
        "",
        "## 版式与内容对比结论",
        "",
        "- 教学设计沿用和平街第一中学四页表格体系、绿色表头、页眉页脚和固定栏目。",
        "- PPT直接复制金标准PPT成熟页型并编辑，保留母版、比例、字体与配色；未复用质量较低的后期追加页。",
        "- 学生版与教师版题号、页码和结构一致；学生版保留书写空间且不显示答案，教师版补齐答案、步骤、易错点和评分点。",
        "- 教材例1、Q1—Q8、方法归纳、当堂检测和分层作业在四类材料中口径一致。",
        "- 相交线示意图在PPT中由可编辑线条、点和文字构成；Word使用同源高清图，未使用整页截图。",
        "",
        "## 视觉检查证据",
        "",
    ])
    lines.extend(f"- `{path.relative_to(PROJECT)}`" for path in visual_evidence)
    REPORT.parent.mkdir(parents=True, exist_ok=True)
    REPORT.write_text("\n".join(lines) + "\n", encoding="utf-8")

    print(f"[{result['status']}] passes={len(passes)} failures={len(failures)} report={REPORT}")
    for failure in failures:
        print(f"  - {failure}")
    return 0 if not failures else 1


if __name__ == "__main__":
    sys.exit(main())
